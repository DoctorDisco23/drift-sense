"""
Script to generate a premium, high-aesthetic PowerPoint presentation (.pptx)
for the Drift-Sense project.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ---------------------------------------------------------
# CONSTANTS & COLOR PALETTE (Dark Tech Aesthetic)
# ---------------------------------------------------------
BG_COLOR = RGBColor(15, 23, 42)        # Deep Navy / Dark Slate (#0F172A)
CARD_BG = RGBColor(30, 41, 59)        # Charcoal Navy (#1E293B)
CARD_BORDER = RGBColor(51, 65, 85)    # Border Slate (#334155)

TEXT_LIGHT = RGBColor(248, 250, 252)  # Crisp White (#F8FAFC)
TEXT_MUTED = RGBColor(148, 163, 184)  # Muted Blue Gray (#94A3B8)

ACCENT_CYAN = RGBColor(14, 165, 233)  # Electric Cyan (#0EA5E9)
ACCENT_PURPLE = RGBColor(139, 92, 246)# Violet (#8B5CF6)
ACCENT_EMERALD = RGBColor(16, 185, 129)# Emerald (#10B981)
ACCENT_AMBER = RGBColor(245, 158, 11) # Amber (#F59E0B)

FONT_HEADING = "Segoe UI"
FONT_BODY = "Calibri"


def create_blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    
    # Add dark background shape
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_COLOR
    bg.line.fill.background()
    return slide


def add_header(slide, title_text, category_text="DRIFT-SENSE / SEMICONDUCTOR ALGORITHMIC TOOLING"):
    # Category Tag
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.4))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.name = FONT_HEADING
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = ACCENT_CYAN
    
    # Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = FONT_HEADING
    p_title.font.size = Pt(26)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_LIGHT


def add_card(slide, left, top, width, height, title, items, accent_color=ACCENT_CYAN):
    # Card Background Shape
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = CARD_BORDER
    card.line.width = Pt(1.5)
    
    # Accent top border strip
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.08))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent_color
    accent_bar.line.fill.background()
    
    # Text Frame inside card
    tx_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), height - Inches(0.3))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    # Title
    p_title = tf.paragraphs[0]
    p_title.text = title
    p_title.font.name = FONT_HEADING
    p_title.font.size = Pt(16)
    p_title.font.bold = True
    p_title.font.color.rgb = accent_color
    p_title.space_after = Pt(10)
    
    # Bullet Items
    for item in items:
        p = tf.add_paragraph()
        if isinstance(item, tuple):
            label, body = item
            run1 = p.add_run()
            run1.text = f"• {label}: "
            run1.font.bold = True
            run1.font.color.rgb = TEXT_LIGHT
            run1.font.size = Pt(12)
            
            run2 = p.add_run()
            run2.text = body
            run2.font.color.rgb = TEXT_MUTED
            run2.font.size = Pt(12)
        else:
            run = p.add_run()
            run.text = f"• {item}"
            run.font.color.rgb = TEXT_MUTED
            run.font.size = Pt(12)
        p.space_after = Pt(6)


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ---------------------------------------------------------
    # SLIDE 1: Title Slide
    # ---------------------------------------------------------
    slide1 = create_blank_slide(prs)
    
    # Decorative Glow Pill
    glow = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(4.2), Inches(0.4))
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(14, 165, 233)
    glow.line.fill.background()
    p_glow = glow.text_frame.paragraphs[0]
    p_glow.text = "SEMICONDUCTOR TOOLING ALGORITHMS"
    p_glow.font.name = FONT_HEADING
    p_glow.font.size = Pt(11)
    p_glow.font.bold = True
    p_glow.font.color.rgb = BG_COLOR
    p_glow.alignment = PP_ALIGN.CENTER
    
    # Main Title
    t_box = slide1.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(1.8))
    tf = t_box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "Drift-Sense"
    p1.font.name = FONT_HEADING
    p1.font.size = Pt(54)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_LIGHT
    
    p2 = tf.add_paragraph()
    p2.text = "Navigation-Error Recovery & Sub-Pixel Alignment for Wafer Inspection Tools"
    p2.font.name = FONT_HEADING
    p2.font.size = Pt(22)
    p2.font.color.rgb = ACCENT_CYAN
    p2.space_before = Pt(8)
    
    # Subtitle / Details Card
    card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.5), Inches(11.733), Inches(2.2))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = CARD_BORDER
    
    tb_card = slide1.shapes.add_textbox(Inches(1.1), Inches(4.7), Inches(11.1), Inches(1.8))
    tf_c = tb_card.text_frame
    tf_c.word_wrap = True
    
    p_c1 = tf_c.paragraphs[0]
    p_c1.text = "PROJECT CORE DELIVERABLES & TEAM ROLES"
    p_c1.font.name = FONT_HEADING
    p_c1.font.size = Pt(14)
    p_c1.font.bold = True
    p_c1.font.color.rgb = ACCENT_PURPLE
    p_c1.space_after = Pt(10)
    
    p_c2 = tf_c.add_paragraph()
    p_c2.text = "• Member B (Algorithmic Lead): Preprocessing (CLAHE), Pyramidal Multi-Scale Search, NMS Peak Extraction, Periodicity Tie-Breaking, Sub-Pixel Parabolic Fitting, ORB+RANSAC AI Fallback, and Technical Documentation."
    p_c2.font.size = Pt(13)
    p_c2.font.color.rgb = TEXT_LIGHT
    p_c2.space_after = Pt(6)
    
    p_c3 = tf_c.add_paragraph()
    p_c3.text = "• Shubh (Data & Evaluation): Synthetic Ground-Truth Data Pipeline, Metrics CSV Generation, Visualization Overlays & Repository Setup."
    p_c3.font.size = Pt(13)
    p_c3.font.color.rgb = TEXT_MUTED

    # ---------------------------------------------------------
    # SLIDE 2: Problem Context & Industry Challenge
    # ---------------------------------------------------------
    slide2 = create_blank_slide(prs)
    add_header(slide2, "Industry Challenge: Stage Drift in Wafer Inspection")
    
    add_card(slide2, Inches(0.8), Inches(1.8), Inches(3.64), Inches(5.0), 
             "1. Mechanical Drift", [
                 ("Stage Hysteresis", "Nanometer-scale mechanical vibration & stage motor drift."),
                 ("Thermal Shifts", "Temperature fluctuations cause micro-shifts across optical fields."),
                 ("Position Drift", "Optical tool fails to land at exact nominal site coordinates.")
             ], ACCENT_CYAN)
             
    add_card(slide2, Inches(4.84), Inches(1.8), Inches(3.64), Inches(5.0), 
             "2. Optical Ambiguity", [
                 ("High Periodicity", "Semiconductor memory dies contain thousands of identical cell structures."),
                 ("Low Contrast", "Silicon oxide/nitride layers produce low optical contrast."),
                 ("Illumination Noise", "Sensor gain differences and dark currents degrade standard matchers.")
             ], ACCENT_PURPLE)
             
    add_card(slide2, Inches(8.88), Inches(1.8), Inches(3.64), Inches(5.0), 
             "3. Technical Solution", [
                 ("Drift-Sense Engine", "Deterministic, sub-pixel matching algorithm requiring 0ms training time."),
                 ("Nanometer Precision", "Sub-0.03 pixel localization accuracy (~15 nm physical resolution)."),
                 ("Scale Resilient", "Pyramidal search handles ±15% optical magnification shifts.")
             ], ACCENT_EMERALD)

    # ---------------------------------------------------------
    # SLIDE 3: System Pipeline Architecture
    # ---------------------------------------------------------
    slide3 = create_blank_slide(prs)
    add_header(slide3, "System Architecture & Processing Pipeline")
    
    # 4 horizontal process cards
    card_w = Inches(2.78)
    card_h = Inches(5.0)
    
    add_card(slide3, Inches(0.8), Inches(1.8), card_w, card_h,
             "STEP 1: PREPROCESS", [
                 ("Input Capture", "Search S (512x512) & Ref R (320x320)"),
                 ("Adaptive CLAHE", "Histogram equalization over 8x8 tiles"),
                 ("Gaussian Blur", "High-frequency noise reduction"),
                 ("Sobel Edges", "Gradient magnitude extraction")
             ], ACCENT_CYAN)
             
    add_card(slide3, Inches(3.78), Inches(1.8), card_w, card_h,
             "STEP 2: MULTI-SCALE", [
                 ("Pyramidal Search", "Scale range: s in [0.85, 1.15]"),
                 ("Bicubic Scaling", "High-quality template interpolation"),
                 ("NCC Map", "Normalized Cross-Correlation evaluation"),
                 ("Response Tensor", "Multi-resolution correlation space")
             ], ACCENT_PURPLE)

    add_card(slide3, Inches(6.76), Inches(1.8), card_w, card_h,
             "STEP 3: NMS & TIE-BREAK", [
                 ("NMS Filtering", "Local dilation filter extracts top-K peaks"),
                 ("Periodicity Check", "Identifies competing grid peaks"),
                 ("Stage Prior", "Center distance penalty score"),
                 ("High-Res Rescore", "Pixel difference verification")
             ], ACCENT_AMBER)

    add_card(slide3, Inches(9.74), Inches(1.8), card_w, card_h,
             "STEP 4: SUB-PIXEL & AI", [
                 ("2D Parabola Fit", "Sub-pixel parabolic quadratic curve fit"),
                 ("Fractional Offset", "Returns dx, dy sub-pixel shifts"),
                 ("AI Fallback", "ORB Keypoints + RANSAC Homography"),
                 ("Final Output", "Precise (x, y) & Confidence score")
             ], ACCENT_EMERALD)

    # ---------------------------------------------------------
    # SLIDE 4: Algorithmic Deep-Dive (Member B Core)
    # ---------------------------------------------------------
    slide4 = create_blank_slide(prs)
    add_header(slide4, "Member B Algorithmic Core: Preprocessing & Pyramids")
    
    add_card(slide4, Inches(0.8), Inches(1.8), Inches(5.66), Inches(5.0),
             "Adaptive Intensity Preprocessing (src/preprocess.py)", [
                 ("CLAHE Contrast Equalization", "Enhances local contrast in low-reflectance silicon oxide/metal layers without blowing out highlights."),
                 ("Gaussian Denoising (3x3)", "Attenuates high-frequency optical sensor shot noise."),
                 ("Sobel Gradient Magnitude", "Extracts spatial edge structures to ensure robust matching even under non-linear lighting variation."),
                 ("Z-Score Calibration", "Normalizes intensity variations for zero mean and unit variance.")
             ], ACCENT_CYAN)

    add_card(slide4, Inches(6.86), Inches(1.8), Inches(5.66), Inches(5.0),
             "Pyramidal Multi-Scale Search Space (src/matcher.py)", [
                 ("Scale Invariance", "Evaluates template scale hypotheses across s in {s_base * 0.95, s_base, s_base * 1.05}."),
                 ("Bicubic & Area Resampling", "Uses INTER_AREA for downscaling and INTER_CUBIC for upscaling to prevent aliasing."),
                 ("Normalized Cross-Correlation (NCC)", "Computes scale-normalized inner product invariant to linear gain shifts."),
                 ("0ms Training Latency", "Fully deterministic signal processing engine with zero GPU training overhead.")
             ], ACCENT_PURPLE)

    # ---------------------------------------------------------
    # SLIDE 5: Periodicity & Sub-Pixel Accuracy (Member B Core)
    # ---------------------------------------------------------
    slide5 = create_blank_slide(prs)
    add_header(slide5, "Periodicity Tie-Breaking & Sub-Pixel Parabolic Fit")
    
    add_card(slide5, Inches(0.8), Inches(1.8), Inches(5.66), Inches(5.0),
             "Periodicity & Stage Prior Tie-Breaking (src/peaks.py)", [
                 ("Non-Maximum Suppression (NMS)", "Extracts local peak candidates using 2D dilation morphology."),
                 ("Grid Ambiguity Detection", "Detects when multiple candidate peaks exhibit correlation within 5% of global maximum."),
                 ("Physical Stage Prior", "Applies physical stage drift prior penalty favoring peaks closer to search ROI center."),
                 ("Composite Fitness Score", "Fitness = 0.8 * NormScore + 0.2 * (1 - NormDistPenalty).")
             ], ACCENT_AMBER)

    add_card(slide5, Inches(6.86), Inches(1.8), Inches(5.66), Inches(5.0),
             "Sub-Pixel Quadratic Surface Fit (src/peaks.py)", [
                 ("Discretization Limit", "Camera sensor pixels discritize position into 1px steps (~500nm)."),
                 ("2D Parabolic Interpolation", "Fits continuous quadratic surface through 3x3 correlation neighborhood around integer peak."),
                 ("Fractional Shift Equation", "dx = (Left - Right) / (2 * (Left - 2*Center + Right))."),
                 ("Extreme Precision", "Achieves nanometer-level spatial localization (< 0.03 px error).")
             ], ACCENT_EMERALD)

    # ---------------------------------------------------------
    # SLIDE 6: AI Fallback & Resilience
    # ---------------------------------------------------------
    slide6 = create_blank_slide(prs)
    add_header(slide6, "AI Feature Matching Fallback & Resilience")
    
    add_card(slide6, Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0),
             "Zero-Shot AI Feature Matching Engine (ORB + RANSAC)", [
                 ("Dynamic Trigger", "Engages automatically when normalized cross-correlation confidence drops below 0.40 under extreme distortion."),
                 ("ORB Keypoint Extraction", "Detects 1,000 scale-invariant oriented FAST keypoints and binary BRIEF descriptors."),
                 ("FLANN / Brute-Force Matcher", "Computes Hamming distance between reference and search keypoint descriptors."),
                 ("RANSAC Homography Estimation", "Filters geometric outliers using 4-point perspective RANSAC (threshold = 5.0 px)."),
                 ("Perspective Projection", "Transforms template center coordinate through homography matrix H to locate target with high confidence.")
             ], ACCENT_PURPLE)

    # ---------------------------------------------------------
    # SLIDE 7: Empirical Benchmark Results
    # ---------------------------------------------------------
    slide7 = create_blank_slide(prs)
    add_header(slide7, "Empirical Benchmark & Performance Verification")
    
    # 4 Metric Cards Across Top
    m_w = Inches(2.78)
    m_h = Inches(1.8)
    
    add_card(slide7, Inches(0.8), Inches(1.8), m_w, m_h, "MEDIAN ERROR", [("0.029 px", "Sub-15nm physical resolution")], ACCENT_EMERALD)
    add_card(slide7, Inches(3.78), Inches(1.8), m_w, m_h, "HIT RATE (< 1px)", [("100.0%", "On unique wafer pattern sites")], ACCENT_CYAN)
    add_card(slide7, Inches(6.76), Inches(1.8), m_w, m_h, "HIT RATE (< 5px)", [("80.0%", "Overall across all 40 test cases")], ACCENT_PURPLE)
    add_card(slide7, Inches(9.74), Inches(1.8), m_w, m_h, "TRAINING LATENCY", [("0 ms", "Deterministic zero-shot inference")], ACCENT_AMBER)
    
    # Lower Table Summary Card
    add_card(slide7, Inches(0.8), Inches(3.8), Inches(11.733), Inches(3.0),
             "Benchmark Execution Results Summary (data/synthetic/ground_truth.csv)", [
                 ("Sample #000 (Easy Unique Pattern)", "True: (111.0, 221.0) | Pred: (110.9601, 221.0228) | Error: 0.046 px | Score: 0.8553"),
                 ("Sample #015 (Easy Unique Pattern)", "True: (207.0, 373.0) | Pred: (207.0207, 373.0123) | Error: 0.024 px | Score: 0.8508"),
                 ("Sample #036 (Hard Periodic Pattern)", "True: (432.0, 112.0) | Pred: (431.9930, 111.9910) | Error: 0.011 px | Score: 0.8548"),
                 ("Overall Dataset Summary", "40 Test Cases | Median Error: 0.029 px | Average Match Confidence: 0.845")
             ], ACCENT_CYAN)

    # ---------------------------------------------------------
    # SLIDE 8: Technical Demo & Q&A Walkthrough
    # ---------------------------------------------------------
    slide8 = create_blank_slide(prs)
    add_header(slide8, "Technical Presentation Walkthrough & Live Demo")
    
    add_card(slide8, Inches(0.8), Inches(1.8), Inches(5.66), Inches(5.0),
             "Demo Video Structure (3-Minute Segment)", [
                 ("Segment 1 (0:00 - 0:45)", "Introduction & Wafer Inspection Stage Drift Problem Context."),
                 ("Segment 2 (0:45 - 1:30)", "Preprocessing & Pyramidal Multi-Scale Search Walkthrough."),
                 ("Segment 3 (1:30 - 2:20)", "NMS Peak Detection, Periodicity Tie-Breaking & Sub-Pixel Parabolic Fit."),
                 ("Segment 4 (2:20 - 3:00)", "ORB+RANSAC AI Fallback & Benchmark Metrics Review.")
             ], ACCENT_CYAN)

    add_card(slide8, Inches(6.86), Inches(1.8), Inches(5.66), Inches(5.0),
             "Technical Q&A Preparation", [
                 ("Q1: Why NCC over Standard Matching?", "NCC is invariant to linear intensity scaling (I -> aI + b), preserving accuracy under optical lighting shifts."),
                 ("Q2: How does Sub-Pixel Fitting work?", "Extracts 3x3 correlation matrix around integer peak and computes 1D quadratic offsets along X and Y axes."),
                 ("Q3: How are Non-Linear Distortions handled?", "ORB keypoints + RANSAC homography project template center coordinates through affine/perspective space.")
             ], ACCENT_PURPLE)

    # ---------------------------------------------------------
    # SLIDE 9: Conclusion & Future Roadmap
    # ---------------------------------------------------------
    slide9 = create_blank_slide(prs)
    add_header(slide9, "Conclusion & Future Hardware Roadmap")
    
    add_card(slide9, Inches(0.8), Inches(1.8), Inches(3.64), Inches(5.0),
             "1. Production Ready", [
                 ("High Accuracy", "Median positioning error of 0.029 pixels."),
                 ("Sub-Pixel Resolution", "Resolves sub-15nm displacement."),
                 ("Zero Training Cost", "Instant deployment without model weights.")
             ], ACCENT_EMERALD)

    add_card(slide9, Inches(4.84), Inches(1.8), Inches(3.64), Inches(5.0),
             "2. Robust Engineering", [
                 ("Multi-Scale Pyramid", "Handles magnification drift."),
                 ("Periodicity Prior", "Resolves repeating memory arrays."),
                 ("AI Fallback", "Recovers from heavy optical noise.")
             ], ACCENT_CYAN)

    add_card(slide9, Inches(8.88), Inches(1.8), Inches(3.64), Inches(5.0),
             "3. Future Acceleration", [
                 ("FPGA / ASIC IP", "Hardware-accelerated 2D NCC correlation pipeline."),
                 ("Sub-Graph Super-Res", "Deep learning sub-pixel feature enhancement."),
                 ("Multi-Die Alignment", "Simultaneous multi-ROI wafer alignment.")
             ], ACCENT_PURPLE)

    out_path = os.path.join("docs", "drift_sense_presentation.pptx")
    prs.save(out_path)
    print(f"Successfully generated PowerPoint presentation at: {out_path}")


if __name__ == "__main__":
    build_presentation()
